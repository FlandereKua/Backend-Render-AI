import pypdf
import docx
import openpyxl
import pptx
import pandas as pd
from pathlib import Path

def parse_pdf(file_path: Path) -> str:
    """Đọc và trích xuất toàn bộ nội dung văn bản từ file PDF."""
    text = ""
    try:
        reader = pypdf.PdfReader(file_path)
        for page in reader.pages:
            text += page.extract_text() + "\n"
    except Exception as e:
        return f"[Lỗi khi đọc file PDF: {e}]"
    return text

def parse_docx(file_path: Path) -> str:
    """Đọc và trích xuất toàn bộ nội dung văn bản từ file DOCX."""
    text = ""
    try:
        doc = docx.Document(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception as e:
        return f"[Lỗi khi đọc file DOCX: {e}]"
    return text

def parse_xlsx(file_path: Path) -> str:
    """
    Sử dụng Pandas để đọc tất cả các sheet từ file Excel và chuyển thành chuỗi văn bản.
    Đây là cách làm mạnh mẽ và ổn định hơn.
    """
    text = ""
    try:
        # engine='openpyxl' được chỉ định để đảm bảo khả năng tương thích
        xls = pd.ExcelFile(file_path, engine='openpyxl')
        for sheet_name in xls.sheet_names:
            text += f"--- Sheet: {sheet_name} ---\n"
            # Đọc từng sheet vào một DataFrame
            df = pd.read_excel(xls, sheet_name=sheet_name)
            # Chuyển DataFrame thành một chuỗi văn bản dễ đọc
            text += df.to_string() + "\n\n"
    except Exception as e:
        return f"[Lỗi khi đọc file XLSX: {e}]"
    return text

def parse_pptx(file_path: Path) -> str:
    """Đọc và trích xuất toàn bộ nội dung văn bản từ các slide trong file PPTX."""
    text = ""
    try:
        prs = pptx.Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            text += f"--- Slide {i+1} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        return f"[Lỗi khi đọc file PPTX: {e}]"
    return text

def parse_file(file_path: Path) -> str:
    """
    Hàm chính để nhận diện loại file và gọi hàm xử lý tương ứng.
    """
    extension = file_path.suffix.lower()
    
    if extension == ".pdf":
        return parse_pdf(file_path)
    elif extension == ".docx":
        return parse_docx(file_path)
    elif extension == ".xlsx":
        return parse_xlsx(file_path)
    elif extension == ".pptx":
        return parse_pptx(file_path)
    elif extension in ['.txt', '.md', '.py', '.js', '.html', '.css', '.json', '.csv']:
        try:
            return file_path.read_text(encoding="utf-8")
        except Exception:
            # Thử đọc với encoding khác nếu utf-8 thất bại
            try:
                return file_path.read_text(encoding="latin-1")
            except Exception as e:
                return f"[Lỗi khi đọc file văn bản: {e}]"
    else:
        raise ValueError(f"Định dạng file '{extension}' không được hỗ trợ.")